"""
create_table_pptx — build a PPTX presentation with a table (and optional text).

Layouts
-------
- "table_only"   : title + full-slide table (default)
- "text_above"   : title + text block + table below
- "text_left"    : title + text on left half + table on right half

Overflow handling
-----------------
1. Shrink font down to MIN_FONT_PT to fit on one slide.
2. If still too many rows, split into multiple slides (each with its own header row).
"""

from __future__ import annotations

import io
import time
import uuid
from typing import Literal

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ── Constants ────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.4)
TITLE_H = Inches(0.7)
TITLE_TOP = Inches(0.2)
CONTENT_TOP = TITLE_TOP + TITLE_H + Inches(0.15)
CONTENT_H = SLIDE_H - CONTENT_TOP - MARGIN   # space below title

ROW_H = Inches(0.38)
HEADER_FONT = 12
DATA_FONT = 11
MIN_FONT_PT = 7          # smallest readable font before we paginate
MAX_ROWS_HARD = 60       # safety cap per slide even at min font

HEADER_BG = RGBColor(0x2F, 0x54, 0x96)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ALT_BG    = RGBColor(0xDD, 0xE8, 0xF5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

Layout = Literal["table_only", "text_above", "text_left"]

# ── In-memory store ──────────────────────────────────────────────────────────

_store: dict[str, tuple[bytes, float]] = {}


def store_pptx(data: bytes, ttl: int) -> str:
    token = uuid.uuid4().hex
    _store[token] = (data, time.monotonic() + ttl)
    return token


def get_pptx(token: str) -> bytes | None:
    entry = _store.get(token)
    if entry is None:
        return None
    data, expiry = entry
    if time.monotonic() > expiry:
        del _store[token]
        return None
    return data


def evict_expired() -> None:
    now = time.monotonic()
    for k in [k for k, (_, exp) in _store.items() if now > exp]:
        del _store[k]


# ── Helpers ──────────────────────────────────────────────────────────────────

def _add_title(slide, text: str) -> None:
    txb = slide.shapes.add_textbox(MARGIN, TITLE_TOP, SLIDE_W - MARGIN * 2, TITLE_H)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.text = text
    runs = tf.paragraphs[0].runs
    if runs:
        runs[0].font.size = Pt(24)
        runs[0].font.bold = True


def _add_text_box(slide, text: str, left, top, width, height) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)


def _rows_per_slide(available_h, font_pt: float) -> int:
    rh = Inches(font_pt / 72 * 1.8)   # approx row height from font size
    return max(1, int(available_h / rh) - 1)  # -1 for header


def _col_widths(headers: list[str], rows: list[list[str]], total_width) -> list[int]:
    """Distribute total_width proportionally by max char length per column."""
    n_cols = len(headers)
    max_lens = [len(str(h)) for h in headers]
    for row in rows:
        for c in range(n_cols):
            max_lens[c] = max(max_lens[c], len(str(row[c])) if c < len(row) else 0)
    total_chars = sum(max_lens) or 1
    min_w = int(total_width / n_cols * 0.4)  # floor: 40% of even share
    widths = [max(min_w, int(total_width * l / total_chars)) for l in max_lens]
    # Scale down proportionally if total exceeds available width
    total = sum(widths)
    if total > total_width:
        scale = total_width / total
        widths = [max(min_w, int(w * scale)) for w in widths]
    # Last column absorbs rounding remainder, clamped to min_w
    widths[-1] = max(min_w, total_width - sum(widths[:-1]))
    return widths


def _add_table(slide, headers, rows, left, top, width, height, font_pt: float) -> None:
    n_cols = len(headers)
    n_rows = len(rows) + 1

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, int(height)).table

    # Apply proportional column widths
    col_widths = _col_widths(headers, rows, int(width))
    for c, cw in enumerate(col_widths):
        tbl.columns[c].width = cw

    for c, text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = str(text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        runs = cell.text_frame.paragraphs[0].runs
        if runs:
            runs[0].font.bold = True
            runs[0].font.size = Pt(font_pt)
            runs[0].font.color.rgb = HEADER_FG

    for r, row in enumerate(rows):
        bg = ALT_BG if r % 2 == 0 else WHITE
        for c in range(n_cols):
            cell = tbl.cell(r + 1, c)
            cell.text = str(row[c]) if c < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            runs = cell.text_frame.paragraphs[0].runs
            if runs:
                runs[0].font.size = Pt(font_pt)


def _compute_font_and_chunks(
    n_rows: int,
    available_h,
    font_pt: float = DATA_FONT,
) -> tuple[float, list[int]]:
    """
    Return (font_pt, chunk_sizes) where chunk_sizes is the number of data
    rows per slide.  Shrinks font first; paginates if still needed.
    """
    # Try shrinking font
    while font_pt >= MIN_FONT_PT:
        rph = _rows_per_slide(available_h, font_pt)
        if rph >= n_rows:
            return font_pt, [n_rows]
        if font_pt == MIN_FONT_PT:
            break
        font_pt -= 1

    # Font at minimum — paginate
    rph = min(_rows_per_slide(available_h, MIN_FONT_PT), MAX_ROWS_HARD)
    chunks = []
    remaining = n_rows
    while remaining > 0:
        take = min(rph, remaining)
        chunks.append(take)
        remaining -= take
    return MIN_FONT_PT, chunks


# ── Main build function ──────────────────────────────────────────────────────

def build_pptx(
    title: str,
    headers: list[str],
    rows: list[list[str]],
    text: str = "",
    layout: Layout = "table_only",
) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    # Determine geometry based on layout
    if layout == "text_above" and text:
        text_h = Inches(1.4)
        tbl_top = CONTENT_TOP + text_h + Inches(0.15)
        tbl_left = MARGIN
        tbl_w = SLIDE_W - MARGIN * 2
        tbl_h = SLIDE_H - tbl_top - MARGIN
        available_h = tbl_h
    elif layout == "text_left" and text:
        col_w = (SLIDE_W - MARGIN * 3) / 2
        tbl_top = CONTENT_TOP
        tbl_left = MARGIN + col_w + MARGIN
        tbl_w = col_w
        tbl_h = CONTENT_H
        available_h = tbl_h
    else:
        layout = "table_only"
        tbl_top = CONTENT_TOP
        tbl_left = MARGIN
        tbl_w = SLIDE_W - MARGIN * 2
        tbl_h = CONTENT_H
        available_h = tbl_h

    font_pt, chunks = _compute_font_and_chunks(len(rows), available_h)

    offset = 0
    for i, chunk_size in enumerate(chunks):
        slide = prs.slides.add_slide(blank_layout)

        page_title = title if len(chunks) == 1 else f"{title} ({i + 1}/{len(chunks)})"
        _add_title(slide, page_title)

        # Text area (only on first slide for multi-page)
        if text and i == 0:
            if layout == "text_above":
                _add_text_box(slide, text, MARGIN, CONTENT_TOP, SLIDE_W - MARGIN * 2, Inches(1.4))
            elif layout == "text_left":
                col_w = (SLIDE_W - MARGIN * 3) / 2
                _add_text_box(slide, text, MARGIN, CONTENT_TOP, col_w, CONTENT_H)

        chunk_rows = rows[offset: offset + chunk_size]
        offset += chunk_size

        # Recalc row height for this chunk
        n_rows_slide = len(chunk_rows) + 1
        actual_h = min(tbl_h, ROW_H * n_rows_slide)

        _add_table(slide, headers, chunk_rows, tbl_left, tbl_top, tbl_w, actual_h, font_pt)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── OpenAI function-calling schema ───────────────────────────────────────────

TOOL_SCHEMA = {
    "type": "function",
    "function": {
        "name": "create_table_pptx",
        "description": (
            "Generate a PowerPoint (.pptx) file containing a table with optional explanatory text. "
            "Use this when the user wants to summarise, compare, or organise data in tabular form."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Specific slide title shown above the table.",
                },
                "headers": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Column header labels.",
                },
                "rows": {
                    "type": "array",
                    "items": {"type": "array", "items": {"type": "string"}},
                    "description": "Table rows; each inner array must match headers length.",
                },
                "text": {
                    "type": "string",
                    "description": (
                        "Optional explanatory text to accompany the table "
                        "(e.g. summary, key findings, context). Leave empty if not needed."
                    ),
                },
                "layout": {
                    "type": "string",
                    "enum": ["table_only", "text_above", "text_left"],
                    "description": (
                        "Slide layout. Use 'text_above' when text is a short heading/summary; "
                        "'text_left' when text and table are equally important; "
                        "'table_only' when no text is provided."
                    ),
                },
            },
            "required": ["title", "headers", "rows"],
            "additionalProperties": False,
        },
    },
}
