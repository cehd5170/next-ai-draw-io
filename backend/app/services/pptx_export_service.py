"""
PPTX export service — converts draw.io XML to a PowerPoint PPTX file.

Two conversion strategies are attempted in order:
1.  ``drawio2pptx`` library (preferred — preserves more styling).
2.  Custom python-pptx-based conversion (fallback for unsupported diagrams).

The custom conversion pipeline:
- Parses mxGraphModel XML with lxml.
- Pass 1: renders all *vertex* mxCell elements as PPTX auto-shapes.
- Pass 2: renders all *edge* mxCell elements as line connectors.

Shape mapping uses ``MSO_AUTO_SHAPE_TYPE`` constants from python-pptx.
"""

from __future__ import annotations

import logging
import re
from typing import Any

from lxml import etree

logger = logging.getLogger(__name__)

# EMU (English Metric Units) per CSS pixel at 96 dpi.
EMU_PER_PX: int = 9525  # 914400 / 96


# ---------------------------------------------------------------------------
# draw.io shape name → MSO_AUTO_SHAPE_TYPE mapping
# ---------------------------------------------------------------------------

# Populated lazily once python-pptx is imported.
_SHAPE_MAP: dict[str, Any] | None = None


def _get_shape_map() -> dict[str, Any]:
    """Build and cache the draw.io → MSO_AUTO_SHAPE_TYPE mapping."""
    global _SHAPE_MAP
    if _SHAPE_MAP is not None:
        return _SHAPE_MAP

    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # noqa: PLC0415

    _SHAPE_MAP = {
        # draw.io style key        python-pptx enum
        "rectangle":               MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "rounded=1":               MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "ellipse":                 MSO_AUTO_SHAPE_TYPE.OVAL,
        "rhombus":                 MSO_AUTO_SHAPE_TYPE.DIAMOND,
        "triangle":                MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        "hexagon":                 MSO_AUTO_SHAPE_TYPE.HEXAGON,
        "parallelogram":           MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
        "trapezoid":               MSO_AUTO_SHAPE_TYPE.TRAPEZOID,
        "cylinder":                MSO_AUTO_SHAPE_TYPE.CAN,
        "cloud":                   MSO_AUTO_SHAPE_TYPE.CLOUD,
        "star":                    MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
        "swimlane":                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        # Additional common draw.io shapes
        "mxgraph.flowchart.start_1":  MSO_AUTO_SHAPE_TYPE.OVAL,
        "mxgraph.flowchart.terminate": MSO_AUTO_SHAPE_TYPE.OVAL,
        "mxgraph.flowchart.decision":  MSO_AUTO_SHAPE_TYPE.DIAMOND,
        "mxgraph.flowchart.process":   MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "mxgraph.flowchart.data":      MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM,
    }
    return _SHAPE_MAP


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------


@staticmethod
def _parse_style(style_str: str) -> dict[str, str]:
    """
    Convert a draw.io style string ``"key1=val1;key2=val2;flag;"``
    into a dict.  Flags without values are stored as ``{flag: "1"}``.
    """
    result: dict[str, str] = {}
    if not style_str:
        return result
    for token in style_str.split(";"):
        token = token.strip()
        if not token:
            continue
        if "=" in token:
            k, _, v = token.partition("=")
            result[k.strip()] = v.strip()
        else:
            result[token] = "1"
    return result


def _px_to_emu(px: float | str) -> int:
    """Convert CSS pixels to EMU."""
    return int(float(px)) * EMU_PER_PX


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """
    Convert a ``#RRGGBB`` (or ``#RGB``) hex string to an ``(r, g, b)`` tuple.

    Returns ``(0, 0, 0)`` for invalid input so callers can use the result
    without extra error handling.
    """
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    if len(hex_color) != 6:
        return (0, 0, 0)
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return (r, g, b)
    except ValueError:
        return (0, 0, 0)


def _get_geometry(cell: etree._Element) -> tuple[float, float, float, float]:
    """
    Extract (x, y, width, height) from an mxCell's mxGeometry child.

    Returns ``(0.0, 0.0, 100.0, 60.0)`` as a safe default if geometry is
    absent or has missing attributes.
    """
    geo = cell.find("mxGeometry")
    if geo is None:
        return (0.0, 0.0, 100.0, 60.0)
    return (
        float(geo.get("x", 0)),
        float(geo.get("y", 0)),
        float(geo.get("width", 100)),
        float(geo.get("height", 60)),
    )


def _determine_shape_type(style: dict[str, str]) -> Any:
    """
    Map a draw.io style dict to an ``MSO_AUTO_SHAPE_TYPE`` constant.

    Falls back to ``RECTANGLE`` when no recognisable shape token is found.
    """
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # noqa: PLC0415

    shape_map = _get_shape_map()

    # Check the ``shape=`` value first (named shapes).
    shape_val = style.get("shape", "")
    if shape_val in shape_map:
        return shape_map[shape_val]

    # Check for implicit shape flags (order matters — most specific first).
    if style.get("ellipse") == "1":
        return MSO_AUTO_SHAPE_TYPE.OVAL
    if style.get("rhombus") == "1":
        return MSO_AUTO_SHAPE_TYPE.DIAMOND
    if style.get("triangle") == "1":
        return MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE

    # Rounded rectangle detection.
    rounded = style.get("rounded", "0")
    if rounded not in ("0", "", "false"):
        return MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE

    return MSO_AUTO_SHAPE_TYPE.RECTANGLE


# ---------------------------------------------------------------------------
# PptxExportService
# ---------------------------------------------------------------------------


class PptxExportService:
    """Convert draw.io mxCell XML to PowerPoint PPTX format."""

    EMU_PER_PX: int = EMU_PER_PX

    # Static accessors exposed as instance methods for API compatibility.
    parse_style = staticmethod(_parse_style)
    px_to_emu = staticmethod(_px_to_emu)
    hex_to_rgb = staticmethod(_hex_to_rgb)

    # ------------------------------------------------------------------
    # Public interface
    # ------------------------------------------------------------------

    async def export(self, drawio_xml: str, filename: str = "diagram.pptx") -> bytes:
        """
        Convert *drawio_xml* to PPTX bytes.

        Attempts ``drawio2pptx`` first; falls back to the custom conversion
        pipeline if the library is unavailable or raises.

        Parameters
        ----------
        drawio_xml:
            Full draw.io mxGraphModel XML (or bare mxCell fragment — a
            wrapper is added automatically if missing).
        filename:
            Desired output filename (used only for logging).

        Returns
        -------
        bytes
            Raw PPTX binary.
        """
        import io  # noqa: PLC0415

        if not drawio_xml or not drawio_xml.strip():
            raise ValueError("Cannot export an empty diagram.")

        # Ensure we have a full mxGraphModel document.
        xml = self._ensure_mxgraph_wrapper(drawio_xml)

        # Strategy 1: drawio2pptx library.
        try:
            pptx_bytes = await self._export_via_drawio2pptx(xml, filename)
            logger.debug("PPTX exported via drawio2pptx (%d bytes)", len(pptx_bytes))
            return pptx_bytes
        except Exception as exc:  # noqa: BLE001
            logger.warning(
                "drawio2pptx export failed (%s); falling back to custom conversion",
                exc,
            )

        # Strategy 2: custom python-pptx conversion.
        try:
            root_el = etree.fromstring(xml.encode())
        except etree.XMLSyntaxError as exc:
            raise ValueError(f"Invalid draw.io XML: {exc}") from exc

        prs = self._custom_convert(root_el)

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()
        logger.debug(
            "PPTX exported via custom conversion (%d bytes)", len(pptx_bytes)
        )
        return pptx_bytes

    # ------------------------------------------------------------------
    # drawio2pptx bridge
    # ------------------------------------------------------------------

    async def _export_via_drawio2pptx(self, xml: str, filename: str) -> bytes:
        """Write *xml* to a temp file, run drawio2pptx, read back bytes."""
        import tempfile  # noqa: PLC0415
        from pathlib import Path  # noqa: PLC0415

        from drawio2pptx import convert  # type: ignore[import]  # noqa: PLC0415

        with tempfile.TemporaryDirectory() as tmpdir:
            tmp = Path(tmpdir)
            src = tmp / "diagram.drawio"
            dst = tmp / filename

            src.write_text(xml, encoding="utf-8")
            convert(str(src), str(dst))

            if not dst.exists():
                raise RuntimeError("drawio2pptx produced no output file")

            return dst.read_bytes()

    # ------------------------------------------------------------------
    # Custom conversion pipeline
    # ------------------------------------------------------------------

    def _custom_convert(self, xml_root: etree._Element) -> Any:
        """
        Custom draw.io → python-pptx conversion.

        1.  Parses the mxGraphModel dimensions for slide sizing.
        2.  Pass 1: renders vertex mxCell elements as auto-shapes.
        3.  Pass 2: renders edge mxCell elements as line connectors.
        """
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Emu  # noqa: PLC0415

        prs = Presentation()

        # Determine slide dimensions from mxGraphModel attributes if present.
        page_width_px = float(xml_root.get("pageWidth", 1169))   # A4 landscape ≈
        page_height_px = float(xml_root.get("pageHeight", 827))  # 1169×827 pt

        prs.slide_width = Emu(_px_to_emu(page_width_px))
        prs.slide_height = Emu(_px_to_emu(page_height_px))

        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Build a lookup from cell id → cell element for edge rendering.
        shape_lookup: dict[str, etree._Element] = {}

        cells = xml_root.findall(".//mxCell")

        # Pass 1: vertices.
        for cell in cells:
            if cell.get("vertex") != "1":
                continue
            # Skip root cells (id=0, id=1).
            if cell.get("id") in ("0", "1"):
                continue
            style = _parse_style(cell.get("style", ""))
            self._add_vertex(slide, cell, style)
            shape_lookup[cell.get("id", "")] = cell

        # Pass 2: edges.
        for cell in cells:
            if cell.get("edge") != "1":
                continue
            style = _parse_style(cell.get("style", ""))
            self._add_edge(slide, cell, style, shape_lookup)

        return prs

    # ------------------------------------------------------------------
    # Shape renderers
    # ------------------------------------------------------------------

    def _add_vertex(self, slide: Any, cell: etree._Element, style: dict[str, str]) -> None:
        """Add a shape to *slide* from an mxCell vertex element."""
        from pptx.dml.color import RGBColor  # noqa: PLC0415
        from pptx.util import Emu, Pt  # noqa: PLC0415

        x, y, w, h = _get_geometry(cell)
        label = cell.get("value", "")

        shape_type = _determine_shape_type(style)

        left = Emu(_px_to_emu(x))
        top = Emu(_px_to_emu(y))
        width = Emu(_px_to_emu(w))
        height = Emu(_px_to_emu(h))

        try:
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        except Exception as exc:  # noqa: BLE001
            logger.debug("Could not add shape type %s: %s", shape_type, exc)
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # noqa: PLC0415
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )

        # Fill colour.
        fill_color_str = style.get("fillColor", "")
        if fill_color_str and fill_color_str.lower() not in ("none", "default"):
            rgb = _hex_to_rgb(fill_color_str)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)
        elif fill_color_str.lower() == "none":
            shape.fill.background()

        # Stroke colour.
        stroke_color_str = style.get("strokeColor", "")
        if stroke_color_str and stroke_color_str.lower() not in ("none", "default"):
            rgb = _hex_to_rgb(stroke_color_str)
            shape.line.color.rgb = RGBColor(*rgb)
        elif stroke_color_str.lower() == "none":
            shape.line.fill.background()

        # Label text.
        if label:
            # Strip HTML tags for plain-text rendering.
            clean_label = re.sub(r"<[^>]+>", "", label).strip()
            if clean_label:
                try:
                    tf = shape.text_frame
                    tf.text = clean_label
                    font_size_str = style.get("fontSize", "11")
                    try:
                        tf.paragraphs[0].runs[0].font.size = Pt(float(font_size_str))
                    except (IndexError, ValueError):
                        pass
                except Exception:  # noqa: BLE001
                    pass

    def _add_edge(
        self,
        slide: Any,
        cell: etree._Element,
        style: dict[str, str],
        shape_lookup: dict[str, etree._Element],
    ) -> None:
        """Add a line connector to *slide* from an mxCell edge element."""
        from pptx.dml.color import RGBColor  # noqa: PLC0415
        from pptx.util import Emu  # noqa: PLC0415

        # Determine start/end points.  If source/target cells exist, use
        # their centre points; otherwise fall back to mxPoint geometry.
        def _centre(c: etree._Element) -> tuple[float, float]:
            cx, cy, cw, ch = _get_geometry(c)
            return cx + cw / 2, cy + ch / 2

        source_id = cell.get("source", "")
        target_id = cell.get("target", "")

        if source_id in shape_lookup:
            sx, sy = _centre(shape_lookup[source_id])
        else:
            geo = cell.find("mxGeometry")
            src_pt = geo.find("mxPoint[@as='sourcePoint']") if geo is not None else None
            sx = float(src_pt.get("x", 0)) if src_pt is not None else 0.0
            sy = float(src_pt.get("y", 0)) if src_pt is not None else 0.0

        if target_id in shape_lookup:
            ex, ey = _centre(shape_lookup[target_id])
        else:
            geo = cell.find("mxGeometry")
            tgt_pt = geo.find("mxPoint[@as='targetPoint']") if geo is not None else None
            ex = float(tgt_pt.get("x", 0)) if tgt_pt is not None else 100.0
            ey = float(tgt_pt.get("y", 0)) if tgt_pt is not None else 100.0

        # Draw as a straight-line connector.
        try:
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                Emu(_px_to_emu(sx)),
                Emu(_px_to_emu(sy)),
                Emu(_px_to_emu(ex)),
                Emu(_px_to_emu(ey)),
            )

            # Stroke colour.
            stroke_color_str = style.get("strokeColor", "#000000")
            if stroke_color_str and stroke_color_str.lower() not in ("none", "default"):
                rgb = _hex_to_rgb(stroke_color_str)
                connector.line.color.rgb = RGBColor(*rgb)

            # Stroke width.
            stroke_width_str = style.get("strokeWidth", "1")
            try:
                from pptx.util import Pt  # noqa: PLC0415
                connector.line.width = Pt(float(stroke_width_str))
            except (ValueError, Exception):  # noqa: BLE001
                pass

        except Exception as exc:  # noqa: BLE001
            logger.debug("Could not add connector: %s", exc)

    def _get_shape_type(self, style: dict[str, str]) -> Any:
        """Map a draw.io style dict to an MSO_AUTO_SHAPE_TYPE constant."""
        return _determine_shape_type(style)

    # ------------------------------------------------------------------
    # Helpers (also exposed as static / class methods for testing)
    # ------------------------------------------------------------------

    @staticmethod
    def _ensure_mxgraph_wrapper(xml: str) -> str:
        """
        Return *xml* wrapped in an mxGraphModel document if it isn't already.
        """
        stripped = xml.strip()
        if stripped.startswith("<mxGraphModel"):
            return stripped
        if stripped.startswith("<mxfile"):
            return stripped
        # Bare mxCell fragment — wrap it.
        return (
            '<mxGraphModel>\n'
            '  <root>\n'
            '    <mxCell id="0" />\n'
            '    <mxCell id="1" parent="0" />\n'
            f'    {stripped}\n'
            '  </root>\n'
            '</mxGraphModel>'
        )
